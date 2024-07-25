from fastapi import UploadFile, HTTPException
from io import BytesIO
import pdfplumber
import PyPDF2
from docx import Document
from pptx import Presentation


"""function that read file"""
async def extract_text(file: UploadFile):
    # extract file extention
    file_ext = file.filename.split('.')[-1].lower()

    # extract the file content
    file_content = BytesIO(await file.read())

    extracted_text = ""

    # if the file extension is pdf, extract pdf text
    if file_ext == 'pdf':
        with pdfplumber.open(path_or_fp=file_content) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    extracted_text += page_text.replace('\n', ' ')
        

    # if the file extension is docx
    elif file_ext == 'docx':
        doc = Document(docx=file_content)
        for para in doc.paragraphs:
            para_text = para.text
            if para_text:
                extracted_text += para_text + " "
        
    # if the file extension is pptx
    elif file_ext == 'pptx':
        prs = Presentation(pptx=file_content)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    extracted_text += paragraph.text + " "    
    else:
        raise HTTPException(
            status_code=400,
            detail="Unsupported file type"
        )
    return extracted_text   